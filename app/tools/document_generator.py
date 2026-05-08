"""
document_generator.py — Generate formatted documents (PDF, DOCX, XLSX, PPTX, HTML).

Enables agents to produce professional output beyond plain text:
  - PDF reports with headers, tables, and formatting
  - Word documents with styles and structure
  - Excel spreadsheets with data and formulas
  - PowerPoint decks with title + content slides and themes
  - Styled HTML pages (served locally for Signal URL delivery)

Output written to workspace/output/docs/ and optionally served via
a local HTTP endpoint so Signal can deliver a clickable URL.

Uses existing libraries: reportlab (PDF), python-docx (DOCX), openpyxl (XLSX),
python-pptx (PPTX).
"""

import json
import logging
import os
import time
from datetime import datetime, timezone
from pathlib import Path

from app.paths import WORKSPACE_ROOT

logger = logging.getLogger(__name__)

OUTPUT_DIR = WORKSPACE_ROOT / "output" / "docs"
try:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
except OSError:
    # On the host, WORKSPACE_ROOT may not be writable yet (e.g. running tests
    # before app.paths.ensure_dirs()). Tests + callers can still patch
    # ``OUTPUT_DIR`` on the module before invoking the generators.
    pass

# Docker workspace root retained as a string for the host-path translation
# helper below — this is the path layout INSIDE the container, not on disk.
_WORKSPACE_ROOT = "/app/workspace"

def _host_path(docker_path: str) -> str:
    """Translate Docker path to host path for Signal delivery."""
    from app.config import get_settings
    host_ws = get_settings().workspace_host_path
    if host_ws:
        return docker_path.replace(_WORKSPACE_ROOT, host_ws)
    return docker_path

def _generate_filename(prefix: str, ext: str) -> str:
    ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    return f"{prefix}_{ts}.{ext}"

# ── PDF Generation ────────────────────────────────────────────────────────────

def create_pdf(
    title: str,
    content: str,
    sections: list[dict] | None = None,
    tables: list[dict] | None = None,
    author: str = "BotArmy Agent Team",
) -> dict:
    """Create a formatted PDF report.

    Args:
        title: Document title
        content: Main body text (markdown-ish — supports paragraphs)
        sections: Optional list of {"heading": str, "body": str}
        tables: Optional list of {"headers": [...], "rows": [[...]]}
        author: Author name

    Returns: {"path": str, "host_path": str, "filename": str, "pages": int}
    """
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.colors import HexColor
        from reportlab.lib.units import mm
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
        )
        from reportlab.lib.enums import TA_JUSTIFY

        filename = _generate_filename("report", "pdf")
        filepath = str(OUTPUT_DIR / filename)

        styles = getSampleStyleSheet()
        story = []

        # Title
        story.append(Paragraph(title, styles['Title']))
        story.append(Spacer(1, 12))
        story.append(Paragraph(
            f"Generated {datetime.now(timezone.utc).strftime('%B %d, %Y')} | {author}",
            styles['Normal']
        ))
        story.append(Spacer(1, 20))

        # Main content
        for para in content.split("\n\n"):
            para = para.strip()
            if para:
                if para.startswith("# "):
                    story.append(Paragraph(para[2:], styles['Heading1']))
                elif para.startswith("## "):
                    story.append(Paragraph(para[3:], styles['Heading2']))
                elif para.startswith("### "):
                    story.append(Paragraph(para[4:], styles['Heading3']))
                else:
                    story.append(Paragraph(para, styles['Normal']))
                story.append(Spacer(1, 6))

        # Sections
        if sections:
            for sec in sections:
                story.append(Paragraph(sec.get("heading", ""), styles['Heading2']))
                story.append(Spacer(1, 4))
                for p in sec.get("body", "").split("\n\n"):
                    if p.strip():
                        story.append(Paragraph(p.strip(), styles['Normal']))
                        story.append(Spacer(1, 4))

        # Tables
        if tables:
            for tbl in tables:
                headers = tbl.get("headers", [])
                rows = tbl.get("rows", [])
                data = [headers] + rows
                t = Table(data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), HexColor("#3b82f6")),
                    ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#ffffff")),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('GRID', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1),
                     [HexColor("#ffffff"), HexColor("#f8fafc")]),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('LEFTPADDING', (0, 0), (-1, -1), 6),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ]))
                story.append(t)
                story.append(Spacer(1, 12))

        doc = SimpleDocTemplate(filepath, pagesize=A4)
        doc.build(story)

        return {
            "success": True,
            "path": filepath,
            "host_path": _host_path(filepath),
            "filename": filename,
            "pages": doc.page,
        }
    except Exception as e:
        logger.error(f"PDF generation failed: {e}")
        return {"success": False, "error": str(e)[:300]}

# ── DOCX Generation ───────────────────────────────────────────────────────────

def create_docx(
    title: str,
    content: str,
    sections: list[dict] | None = None,
    tables: list[dict] | None = None,
    author: str = "BotArmy Agent Team",
) -> dict:
    """Create a formatted Word document.

    Returns: {"path": str, "host_path": str, "filename": str}
    """
    try:
        from docx import Document
        from docx.shared import Inches, Pt

        filename = _generate_filename("report", "docx")
        filepath = str(OUTPUT_DIR / filename)

        doc = Document()
        doc.core_properties.author = author
        doc.core_properties.title = title

        doc.add_heading(title, level=0)
        doc.add_paragraph(
            f"Generated {datetime.now(timezone.utc).strftime('%B %d, %Y')} | {author}"
        ).italic = True

        # Main content
        for para in content.split("\n\n"):
            para = para.strip()
            if para:
                if para.startswith("# "):
                    doc.add_heading(para[2:], level=1)
                elif para.startswith("## "):
                    doc.add_heading(para[3:], level=2)
                elif para.startswith("### "):
                    doc.add_heading(para[4:], level=3)
                elif para.startswith("- "):
                    for line in para.split("\n"):
                        if line.strip().startswith("- "):
                            doc.add_paragraph(line.strip()[2:], style='List Bullet')
                else:
                    doc.add_paragraph(para)

        # Sections
        if sections:
            for sec in sections:
                doc.add_heading(sec.get("heading", ""), level=2)
                for p in sec.get("body", "").split("\n\n"):
                    if p.strip():
                        doc.add_paragraph(p.strip())

        # Tables
        if tables:
            for tbl in tables:
                headers = tbl.get("headers", [])
                rows = tbl.get("rows", [])
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = 'Light Grid Accent 1'
                for i, h in enumerate(headers):
                    table.rows[0].cells[i].text = str(h)
                for r_idx, row in enumerate(rows):
                    for c_idx, cell in enumerate(row):
                        table.rows[r_idx + 1].cells[c_idx].text = str(cell)

        doc.save(filepath)
        return {
            "success": True,
            "path": filepath,
            "host_path": _host_path(filepath),
            "filename": filename,
        }
    except Exception as e:
        logger.error(f"DOCX generation failed: {e}")
        return {"success": False, "error": str(e)[:300]}

# ── XLSX Generation ───────────────────────────────────────────────────────────

def create_xlsx(
    title: str,
    sheets: list[dict] | None = None,
    data: list[list] | None = None,
    headers: list[str] | None = None,
) -> dict:
    """Create an Excel spreadsheet.

    Args:
        title: Workbook title / first sheet name
        sheets: List of {"name": str, "headers": [...], "rows": [[...]]}
        data: Simple data (if no sheets provided)
        headers: Column headers (if no sheets provided)

    Returns: {"path": str, "host_path": str, "filename": str}
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment

        filename = _generate_filename("data", "xlsx")
        filepath = str(OUTPUT_DIR / filename)

        wb = Workbook()

        if sheets:
            for i, sheet in enumerate(sheets):
                ws = wb.active if i == 0 else wb.create_sheet()
                ws.title = sheet.get("name", f"Sheet{i+1}")
                hdrs = sheet.get("headers", [])
                rows = sheet.get("rows", [])

                # Headers
                for col, h in enumerate(hdrs, 1):
                    cell = ws.cell(row=1, column=col, value=str(h))
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor="3B82F6")
                    cell.alignment = Alignment(horizontal="center")

                # Data
                for r_idx, row in enumerate(rows, 2):
                    for c_idx, val in enumerate(row, 1):
                        ws.cell(row=r_idx, column=c_idx, value=val)

                # Auto-width
                for col in ws.columns:
                    max_len = max(len(str(c.value or "")) for c in col)
                    ws.column_dimensions[col[0].column_letter].width = min(50, max_len + 2)
        else:
            ws = wb.active
            ws.title = title[:31]
            if headers:
                for col, h in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=str(h))
                    cell.font = Font(bold=True)
            if data:
                start_row = 2 if headers else 1
                for r_idx, row in enumerate(data, start_row):
                    for c_idx, val in enumerate(row, 1):
                        ws.cell(row=r_idx, column=c_idx, value=val)

        wb.save(filepath)
        return {
            "success": True,
            "path": filepath,
            "host_path": _host_path(filepath),
            "filename": filename,
        }
    except Exception as e:
        logger.error(f"XLSX generation failed: {e}")
        return {"success": False, "error": str(e)[:300]}

# ── PPTX Generation ───────────────────────────────────────────────────────────

# Theme palettes — same vocabulary the HTML generator uses, but expressed as
# RGB triples so python-pptx can paint shape fills directly.
_PPTX_THEMES = {
    "modern-dark": {
        "bg":     (0x0F, 0x17, 0x2A),
        "text":   (0xE2, 0xE8, 0xF0),
        "accent": (0x3B, 0x82, 0xF6),
        "muted":  (0x94, 0xA3, 0xB8),
    },
    "clean-light": {
        "bg":     (0xFF, 0xFF, 0xFF),
        "text":   (0x1E, 0x29, 0x3B),
        "accent": (0x25, 0x63, 0xEB),
        "muted":  (0x64, 0x74, 0x8B),
    },
    "minimal": {
        "bg":     (0xFA, 0xFA, 0xFA),
        "text":   (0x17, 0x17, 0x17),
        "accent": (0x00, 0x00, 0x00),
        "muted":  (0x73, 0x73, 0x73),
    },
}


def create_pptx(
    title: str,
    slides: list[dict] | None = None,
    *,
    subtitle: str = "",
    style: str = "modern-dark",
    author: str = "BotArmy Agent Team",
) -> dict:
    """Create a PowerPoint deck (.pptx).

    Args:
        title: Deck title — appears on the cover slide.
        slides: List of content slides. Each is a dict with at least one of:
            ``title`` (str), ``body`` (str — paragraphs separated by blank
            lines, lines starting with ``- `` become bullets), ``bullets``
            (list[str]), ``table`` ({"headers": [...], "rows": [[...]]}),
            ``notes`` (str — speaker notes).
        subtitle: Optional subtitle on the cover slide. Falls back to a
            "Generated <date> | <author>" line.
        style: Visual theme — "modern-dark", "clean-light", or "minimal".
            Same palette vocabulary as ``create_html_page``.
        author: Author attribution.

    Returns:
        ``{"success": True, "path": str, "host_path": str, "filename": str,
           "slides": int}`` on success, ``{"success": False, "error": str}``
        on failure.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        theme = _PPTX_THEMES.get(style, _PPTX_THEMES["modern-dark"])

        def rgb(name: str) -> RGBColor:
            return RGBColor(*theme[name])

        filename = _generate_filename("deck", "pptx")
        filepath = str(OUTPUT_DIR / filename)

        prs = Presentation()
        # Standard 16:9 widescreen — matches Google Slides / Keynote default.
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── Cover slide (blank layout so we control everything) ────────
        cover = prs.slides.add_slide(prs.slide_layouts[6])
        _paint_background(cover, prs, rgb("bg"))
        _accent_bar(cover, prs, rgb("accent"))

        title_box = cover.shapes.add_textbox(
            Inches(0.7), Inches(2.4), Inches(12), Inches(2),
        )
        _set_text(title_box, title, font_size=44, bold=True, color=rgb("text"))

        subtitle_text = subtitle or (
            f"Generated {datetime.now(timezone.utc).strftime('%B %d, %Y')} · {author}"
        )
        sub_box = cover.shapes.add_textbox(
            Inches(0.7), Inches(4.5), Inches(12), Inches(1),
        )
        _set_text(sub_box, subtitle_text, font_size=20, color=rgb("muted"))

        # ── Content slides ─────────────────────────────────────────────
        slide_specs = list(slides or [])
        for spec in slide_specs:
            _add_content_slide(prs, spec, theme)

        prs.save(filepath)
        return {
            "success": True,
            "path": filepath,
            "host_path": _host_path(filepath),
            "filename": filename,
            "slides": 1 + len(slide_specs),
        }
    except Exception as e:
        logger.error(f"PPTX generation failed: {e}", exc_info=True)
        return {"success": False, "error": str(e)[:300]}


def _add_content_slide(prs, spec: dict, theme: dict) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    def rgb(name: str) -> RGBColor:
        return RGBColor(*theme[name])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_background(slide, prs, rgb("bg"))
    _accent_bar(slide, prs, rgb("accent"))

    # Title row
    title_text = (spec.get("title") or "").strip()
    if title_text:
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5), Inches(12), Inches(1),
        )
        _set_text(title_box, title_text, font_size=32, bold=True, color=rgb("text"))

    # Body — bullets array preferred, otherwise parse markdown-ish body string.
    bullets = spec.get("bullets")
    body_text = spec.get("body") or ""
    body_lines: list[tuple[str, int]] = []  # (text, indent_level)
    if isinstance(bullets, list) and bullets:
        body_lines = [(str(b), 0) for b in bullets if str(b).strip()]
    elif body_text:
        for para in body_text.split("\n\n"):
            for line in para.splitlines():
                stripped = line.lstrip()
                if not stripped:
                    continue
                if stripped.startswith("- "):
                    indent = (len(line) - len(stripped)) // 2
                    body_lines.append((stripped[2:], min(indent, 3)))
                else:
                    body_lines.append((stripped, 0))

    if body_lines:
        body_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7), Inches(12), Inches(5),
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, (line, level) in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.level = level
            for run in p.runs:
                run.font.size = Pt(20 if level == 0 else 16)
                run.font.color.rgb = rgb("text")

    # Optional table
    table_spec = spec.get("table")
    if isinstance(table_spec, dict):
        headers = table_spec.get("headers") or []
        rows = table_spec.get("rows") or []
        if headers and rows:
            n_rows = 1 + len(rows)
            n_cols = len(headers)
            top_inches = 1.7 + (1.2 if body_lines else 0)
            tbl_shape = slide.shapes.add_table(
                n_rows, n_cols,
                Inches(0.7), Inches(top_inches),
                Inches(12), Inches(min(5, 0.6 * n_rows + 0.6)),
            )
            tbl = tbl_shape.table
            for c, h in enumerate(headers):
                cell = tbl.cell(0, c)
                cell.text = str(h)
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb("accent")
            for r, row in enumerate(rows, start=1):
                for c, val in enumerate(row[:n_cols]):
                    cell = tbl.cell(r, c)
                    cell.text = str(val)
                    for run in cell.text_frame.paragraphs[0].runs:
                        run.font.size = Pt(12)
                        run.font.color.rgb = rgb("text")

    # Speaker notes
    notes_text = spec.get("notes")
    if notes_text:
        slide.notes_slide.notes_text_frame.text = str(notes_text)


def _paint_background(slide, prs, color) -> None:
    """Fill the slide background with a solid colour."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), prs.slide_width, prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()  # no border
    # Send to back so foreground shapes paint over it.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _accent_bar(slide, prs, color) -> None:
    """Thin accent strip across the top — gives every slide a consistent header."""
    from pptx.util import Inches, Emu
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0), prs.slide_width, Inches(0.18),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _set_text(textbox, text: str, *, font_size: int, color, bold: bool = False) -> None:
    """Set text + font on a fresh textbox's first paragraph."""
    from pptx.util import Pt
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


# ── HTML Page Generation (for Signal URL delivery) ────────────────────────────

def create_html_page(
    title: str,
    content: str,
    sections: list[dict] | None = None,
    tables: list[dict] | None = None,
    style: str = "modern-dark",
    author: str = "BotArmy Agent Team",
) -> dict:
    """Create a styled HTML page and return a local URL.

    The page is saved to workspace/output/docs/ and can be served
    via the gateway's /dashboard endpoint pattern or sent as a file attachment.

    Args:
        title: Page title
        content: Main body (supports markdown-ish headers and paragraphs)
        sections: Optional structured sections
        tables: Optional data tables
        style: Visual style ("modern-dark", "clean-light", "minimal")
        author: Attribution

    Returns: {"path": str, "host_path": str, "filename": str, "url": str}
    """
    filename = _generate_filename("page", "html")
    filepath = str(OUTPUT_DIR / filename)

    # Style themes
    themes = {
        "modern-dark": {
            "bg": "#0f172a", "surface": "#1e293b", "text": "#e2e8f0",
            "accent": "#3b82f6", "muted": "#94a3b8", "border": "#334155",
        },
        "clean-light": {
            "bg": "#ffffff", "surface": "#f8fafc", "text": "#1e293b",
            "accent": "#2563eb", "muted": "#64748b", "border": "#e2e8f0",
        },
        "minimal": {
            "bg": "#fafafa", "surface": "#ffffff", "text": "#171717",
            "accent": "#000000", "muted": "#737373", "border": "#e5e5e5",
        },
    }
    t = themes.get(style, themes["modern-dark"])

    # Build HTML content
    body_html = ""
    for para in content.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        if para.startswith("# "):
            body_html += f'<h1>{_esc(para[2:])}</h1>\n'
        elif para.startswith("## "):
            body_html += f'<h2>{_esc(para[3:])}</h2>\n'
        elif para.startswith("### "):
            body_html += f'<h3>{_esc(para[4:])}</h3>\n'
        elif para.startswith("- "):
            items = "\n".join(f"<li>{_esc(l.strip()[2:])}</li>" for l in para.split("\n") if l.strip().startswith("- "))
            body_html += f"<ul>{items}</ul>\n"
        else:
            body_html += f'<p>{_esc(para)}</p>\n'

    if sections:
        for sec in sections:
            body_html += f'<h2>{_esc(sec.get("heading", ""))}</h2>\n'
            for p in sec.get("body", "").split("\n\n"):
                if p.strip():
                    body_html += f'<p>{_esc(p.strip())}</p>\n'

    if tables:
        for tbl in tables:
            headers = tbl.get("headers", [])
            rows = tbl.get("rows", [])
            body_html += '<div class="table-wrap"><table>\n<thead><tr>'
            for h in headers:
                body_html += f'<th>{_esc(str(h))}</th>'
            body_html += '</tr></thead>\n<tbody>\n'
            for row in rows:
                body_html += '<tr>'
                for cell in row:
                    body_html += f'<td>{_esc(str(cell))}</td>'
                body_html += '</tr>\n'
            body_html += '</tbody></table></div>\n'

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{_esc(title)}</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ background: {t['bg']}; color: {t['text']}; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', system-ui, sans-serif; line-height: 1.6; padding: 0; }}
.container {{ max-width: 900px; margin: 0 auto; padding: 40px 24px; }}
header {{ border-bottom: 2px solid {t['accent']}; padding-bottom: 20px; margin-bottom: 30px; }}
h1 {{ font-size: 2rem; font-weight: 700; margin-bottom: 8px; color: {t['text']}; }}
.meta {{ font-size: 0.85rem; color: {t['muted']}; }}
h2 {{ font-size: 1.4rem; font-weight: 600; margin: 28px 0 12px; color: {t['accent']}; }}
h3 {{ font-size: 1.1rem; font-weight: 600; margin: 20px 0 8px; }}
p {{ margin-bottom: 12px; color: {t['text']}; }}
ul {{ margin: 12px 0 12px 24px; }}
li {{ margin-bottom: 4px; }}
.table-wrap {{ overflow-x: auto; margin: 16px 0; }}
table {{ border-collapse: collapse; width: 100%; font-size: 0.9rem; }}
th {{ background: {t['accent']}; color: white; padding: 10px 12px; text-align: left; font-weight: 600; }}
td {{ padding: 8px 12px; border-bottom: 1px solid {t['border']}; }}
tr:nth-child(even) {{ background: {t['surface']}; }}
footer {{ margin-top: 40px; padding-top: 16px; border-top: 1px solid {t['border']}; font-size: 0.75rem; color: {t['muted']}; }}
@media (max-width: 640px) {{ .container {{ padding: 20px 16px; }} h1 {{ font-size: 1.5rem; }} }}
</style>
</head>
<body>
<div class="container">
<header>
<h1>{_esc(title)}</h1>
<div class="meta">{author} &middot; {datetime.now(timezone.utc).strftime('%B %d, %Y')}</div>
</header>
<main>
{body_html}
</main>
<footer>Generated by BotArmy Agent Team &middot; {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}</footer>
</div>
</body>
</html>"""

    from app.safe_io import safe_write
    safe_write(Path(filepath), html)

    # Generate URL (served via gateway if accessible)
    from app.config import get_settings
    s = get_settings()
    url = f"http://localhost:{s.gateway_port}/docs/{filename}"

    return {
        "success": True,
        "path": filepath,
        "host_path": _host_path(filepath),
        "filename": filename,
        "url": url,
    }

def _esc(text: str) -> str:
    """Escape HTML entities."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")

# ── List generated documents ──────────────────────────────────────────────────

def list_documents() -> list[dict]:
    """List all generated documents."""
    docs = []
    for f in sorted(OUTPUT_DIR.iterdir(), reverse=True):
        if f.is_file():
            docs.append({
                "filename": f.name,
                "path": str(f),
                "size_bytes": f.stat().st_size,
                "created": datetime.fromtimestamp(f.stat().st_mtime, tz=timezone.utc).isoformat(),
                "type": f.suffix[1:],
            })
    return docs[:50]

# ── CrewAI tool wrappers ─────────────────────────────────────────────────────

def create_document_tools() -> list:
    """Create CrewAI-compatible document generation tools for agents."""
    try:
        from crewai.tools import tool

        @tool("generate_pdf")
        def pdf_tool(title: str, content: str, sections: str = "") -> str:
            """Generate a formatted PDF document. Provide title, content text,
            and optionally sections as JSON: [{"heading": "...", "body": "..."}].
            Returns the file path."""
            section_list = None
            if sections:
                try:
                    section_list = json.loads(sections)
                except (json.JSONDecodeError, TypeError):
                    pass
            result = create_pdf(title=title, content=content, sections=section_list)
            return result.get("path", result.get("error", "PDF generation failed"))

        @tool("generate_docx")
        def docx_tool(title: str, content: str, sections: str = "") -> str:
            """Generate a Word document (.docx). Provide title, content,
            and optionally sections as JSON. Returns the file path."""
            section_list = None
            if sections:
                try:
                    section_list = json.loads(sections)
                except (json.JSONDecodeError, TypeError):
                    pass
            result = create_docx(title=title, content=content, sections=section_list)
            return result.get("path", result.get("error", "DOCX generation failed"))

        @tool("generate_html_page")
        def html_tool(title: str, content: str, theme: str = "modern-dark") -> str:
            """Generate a styled HTML page. Themes: modern-dark, clean-light, minimal.
            Returns the URL where the page is served."""
            result = create_html_page(title=title, body_html=content, theme=theme)
            return result.get("url", result.get("path", result.get("error", "HTML generation failed")))

        @tool("generate_pptx_report")
        def pptx_tool(
            title: str,
            slides: str,
            subtitle: str = "",
            theme: str = "modern-dark",
        ) -> str:
            """Generate a PowerPoint deck (.pptx) and return the file path.

            Args:
                title: Cover-slide title.
                slides: JSON array of content slides. Each item:
                    {"title": "...", "body": "bullet/text\\n- item\\n- item",
                     "bullets": ["alt to body — one bullet per item"],
                     "table": {"headers": [...], "rows": [[...]]},
                     "notes": "speaker notes"}.
                subtitle: Optional cover subtitle (else auto-stamps date+author).
                theme: "modern-dark", "clean-light", or "minimal".

            Returns the absolute file path so the deck can be sent as a Signal
            attachment via signal_send_attachment.
            """
            slide_list: list[dict] = []
            if slides:
                try:
                    parsed = json.loads(slides)
                    if isinstance(parsed, list):
                        slide_list = [s for s in parsed if isinstance(s, dict)]
                except (json.JSONDecodeError, TypeError):
                    pass
            result = create_pptx(
                title=title, slides=slide_list,
                subtitle=subtitle, style=theme,
            )
            return result.get("path", result.get("error", "PPTX generation failed"))

        return [pdf_tool, docx_tool, html_tool, pptx_tool]
    except Exception:
        return []
