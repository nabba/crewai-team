"""
Document ingestion pipeline — extract text from many formats, chunk it,
and prepare metadata-enriched chunks for vector storage.
"""

import hashlib
import json
import logging
import os
import re
from dataclasses import dataclass, field
from datetime import datetime, timezone
from typing import Optional

from app.knowledge_base import config

logger = logging.getLogger(__name__)

# Max characters to extract from any single document (safety cap).
_MAX_EXTRACT_CHARS = 500_000


# ─────────────────────────────────────────────────────────────────────────────
# Data classes
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class DocumentChunk:
    text: str
    metadata: dict = field(default_factory=dict)

    @property
    def chunk_id(self) -> str:
        """Deterministic ID based on source + chunk index."""
        source = self.metadata.get("source_path", "unknown")
        idx = self.metadata.get("chunk_index", 0)
        raw = f"{source}::chunk::{idx}"
        return hashlib.sha256(raw.encode()).hexdigest()[:16]


@dataclass
class IngestionResult:
    source: str
    format: str
    chunks_created: int
    total_characters: int
    success: bool
    error: str = ""
    document_id: str = ""


# ─────────────────────────────────────────────────────────────────────────────
# Format detection
# ─────────────────────────────────────────────────────────────────────────────

def detect_format(source: str) -> str:
    """Return a normalised format key for the given source path or URL."""
    if source.startswith(("http://", "https://")):
        return "url"
    ext = os.path.splitext(source)[1].lower()
    if ext in config.SUPPORTED_EXTENSIONS:
        return ext
    return ext or "unknown"


# ─────────────────────────────────────────────────────────────────────────────
# Text extractors
# ─────────────────────────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    import pypdf
    reader = pypdf.PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages[:100]):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)
    return "\n\n".join(pages)[:_MAX_EXTRACT_CHARS]


def extract_docx(path: str) -> str:
    import docx
    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)[:_MAX_EXTRACT_CHARS]


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)[:_MAX_EXTRACT_CHARS]


def extract_xlsx(path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for name in wb.sheetnames[:10]:
        ws = wb[name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i > 500:
                rows.append("... (truncated)")
                break
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"=== Sheet: {name} ===\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)[:_MAX_EXTRACT_CHARS]


def extract_csv(path: str) -> str:
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i > 1000:
                rows.append("... (truncated)")
                break
            rows.append(" | ".join(row))
    return "\n".join(rows)[:_MAX_EXTRACT_CHARS]


def extract_text(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()[:_MAX_EXTRACT_CHARS]


def extract_html(path: str) -> str:
    try:
        import html2text
        with open(path, encoding="utf-8", errors="replace") as f:
            raw = f.read()[:_MAX_EXTRACT_CHARS]
        h = html2text.HTML2Text()
        h.ignore_links = False
        h.ignore_images = True
        return h.handle(raw)
    except ImportError:
        from bs4 import BeautifulSoup
        with open(path, encoding="utf-8", errors="replace") as f:
            raw = f.read()[:_MAX_EXTRACT_CHARS]
        soup = BeautifulSoup(raw, "html.parser")
        return soup.get_text(separator="\n")


def extract_json(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)[:_MAX_EXTRACT_CHARS]


def extract_url(url: str) -> str:
    try:
        import trafilatura
        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            raise ValueError(f"Could not download {url}")
        text = trafilatura.extract(downloaded, include_comments=False)
        if not text:
            raise ValueError(f"Could not extract text from {url}")
        return text[:_MAX_EXTRACT_CHARS]
    except ImportError:
        import requests
        resp = requests.get(url, timeout=30, headers={"User-Agent": "BotArmy-KB/1.0"})
        resp.raise_for_status()
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(resp.text, "html.parser")
        return soup.get_text(separator="\n")[:_MAX_EXTRACT_CHARS]


# Extractor registry
EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
    ".txt": extract_text,
    ".md": extract_text,
    ".html": extract_html,
    ".htm": extract_html,
    ".json": extract_json,
}


# ─────────────────────────────────────────────────────────────────────────────
# Chunking
# ─────────────────────────────────────────────────────────────────────────────

def chunk_text(
    text: str,
    chunk_size: int = config.CHUNK_SIZE,
    chunk_overlap: int = config.CHUNK_OVERLAP,
) -> list[str]:
    """Split text into semantically-aware chunks."""
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except ImportError:
        # Fallback: simple chunking if langchain-text-splitters not installed
        return _simple_chunk(text, chunk_size, chunk_overlap)

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=[
            "\n## ",      # Markdown H2
            "\n### ",     # Markdown H3
            "\n\n",       # Paragraph
            "\n",         # Line
            ". ",         # Sentence
            ", ",         # Clause
            " ",          # Word
            "",           # Char fallback
        ],
        keep_separator=True,
        strip_whitespace=True,
    )
    chunks = splitter.split_text(text)
    return [c for c in chunks if len(c.strip()) > 50]


def _simple_chunk(text: str, size: int, overlap: int) -> list[str]:
    """Fallback chunker without langchain dependency."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + size
        chunk = text[start:end].strip()
        if len(chunk) > 50:
            chunks.append(chunk)
        start = end - overlap
    return chunks


# ─────────────────────────────────────────────────────────────────────────────
# Main ingestion function
# ─────────────────────────────────────────────────────────────────────────────

def ingest_document(
    source: str,
    category: str = "general",
    tags: Optional[list[str]] = None,
    chunk_size: int = config.CHUNK_SIZE,
    chunk_overlap: int = config.CHUNK_OVERLAP,
) -> tuple[list[DocumentChunk], IngestionResult]:
    """
    Process a document from any supported source into vector-ready chunks.
    """
    tags = tags or []
    now = datetime.now(timezone.utc).isoformat()

    try:
        fmt = detect_format(source)

        if fmt == "url":
            raw_text = extract_url(source)
            source_name = source
            file_format = "url"
        else:
            if not os.path.exists(source):
                raise FileNotFoundError(f"File not found: {source}")
            extractor = EXTRACTORS.get(fmt)
            if not extractor:
                raise ValueError(f"Unsupported format: {fmt}")
            raw_text = extractor(source)
            source_name = os.path.basename(source)
            file_format = fmt.lstrip(".")

        if not raw_text or not raw_text.strip():
            return [], IngestionResult(
                source=source,
                format=file_format if fmt != "url" else "url",
                chunks_created=0,
                total_characters=0,
                success=False,
                error="No text content could be extracted.",
            )

        # Clean
        raw_text = re.sub(r"\n{3,}", "\n\n", raw_text).strip()

        # Chunk
        text_chunks = chunk_text(raw_text, chunk_size, chunk_overlap)

        # Build metadata-enriched chunks
        doc_id = hashlib.sha256(source.encode()).hexdigest()[:12]
        doc_chunks = []
        for i, text in enumerate(text_chunks):
            metadata = {
                "source": source_name,
                "source_path": source,
                "format": file_format,
                "category": category,
                "tags": json.dumps(tags),
                "chunk_index": i,
                "total_chunks": len(text_chunks),
                "ingested_at": now,
                "char_count": len(text),
            }
            doc_chunks.append(DocumentChunk(text=text, metadata=metadata))

        result = IngestionResult(
            source=source_name,
            format=file_format,
            chunks_created=len(doc_chunks),
            total_characters=len(raw_text),
            success=True,
            document_id=doc_id,
        )

        logger.info(
            f"Ingested '{source_name}': {len(doc_chunks)} chunks, "
            f"{len(raw_text):,} characters"
        )
        return doc_chunks, result

    except Exception as e:
        logger.error(f"Ingestion failed for '{source}': {e}")
        return [], IngestionResult(
            source=source,
            format="unknown",
            chunks_created=0,
            total_characters=0,
            success=False,
            error=str(e),
        )
