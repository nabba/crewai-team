"""
test_document_generator_pptx — round-trip test for the PPTX generator.

Exercises ``app.tools.document_generator.create_pptx`` end-to-end:
  - generates a deck with cover + multiple content slides
  - re-opens the .pptx via python-pptx
  - confirms slide count, title text, bullet text, table cells, and notes
"""
from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def tmp_output_dir(tmp_path, monkeypatch):
    """Redirect document_generator.OUTPUT_DIR to a tmp dir.

    The module decides its OUTPUT_DIR at import time, so we need to patch it
    on the module after import.
    """
    from app.tools import document_generator as dg
    out = tmp_path / "docs"
    out.mkdir()
    monkeypatch.setattr(dg, "OUTPUT_DIR", out, raising=False)
    return out


def test_create_pptx_round_trip(tmp_output_dir: Path):
    from pptx import Presentation
    from app.tools.document_generator import create_pptx

    slides = [
        {
            "title": "What we shipped",
            "bullets": ["Voice mode", "Vision CU cap", "Concierge persona"],
            "notes": "Speaker note for slide 1.",
        },
        {
            "title": "Cost comparison",
            "body": "Cloud is cheap at low volume.\n\n- Local: $0/mo\n- Cloud: ~$3/mo",
            "table": {
                "headers": ["Mode", "STT/mo", "TTS/mo"],
                "rows": [
                    ["Local",  "$0",    "$0"],
                    ["Cloud",  "$0.60", "$2.70"],
                ],
            },
        },
        {
            "title": "Risks",
            "body": "- Mac mini sleep blocks bridge\n- Whisper-large model is 3 GB",
        },
    ]

    result = create_pptx(
        title="Phase-0 review",
        slides=slides,
        subtitle="Internal",
        style="modern-dark",
        author="AndrusAI",
    )

    assert result["success"] is True, result.get("error")
    assert result["slides"] == 4  # cover + 3 content
    path = Path(result["path"])
    assert path.exists() and path.suffix == ".pptx"

    # Re-open and verify content survived the write/read cycle.
    prs = Presentation(str(path))
    assert len(prs.slides) == 4

    # Cover slide: title text appears
    cover_text = _slide_text(prs.slides[0])
    assert "Phase-0 review" in cover_text
    assert "Internal" in cover_text

    # Slide 1 — bullet list
    s1_text = _slide_text(prs.slides[1])
    assert "What we shipped" in s1_text
    assert "Voice mode" in s1_text
    assert "Vision CU cap" in s1_text
    # Speaker notes
    assert "slide 1" in prs.slides[1].notes_slide.notes_text_frame.text.lower()

    # Slide 2 — body + table
    s2 = prs.slides[2]
    s2_text = _slide_text(s2)
    assert "Cost comparison" in s2_text
    assert "Local" in s2_text and "Cloud" in s2_text
    # Table values
    table_shapes = [shape for shape in s2.shapes if shape.has_table]
    assert len(table_shapes) == 1
    tbl = table_shapes[0].table
    assert tbl.cell(0, 0).text == "Mode"
    assert tbl.cell(0, 1).text == "STT/mo"
    assert tbl.cell(1, 0).text == "Local"
    assert tbl.cell(2, 1).text == "$0.60"

    # Slide 3 — body bullets parsed from markdown
    s3_text = _slide_text(prs.slides[3])
    assert "Risks" in s3_text
    assert "Mac mini sleep blocks bridge" in s3_text


def test_create_pptx_empty_slides_still_writes_cover(tmp_output_dir: Path):
    from app.tools.document_generator import create_pptx
    result = create_pptx(title="Just a cover", slides=None, style="clean-light")
    assert result["success"] is True
    assert result["slides"] == 1


def test_create_pptx_unknown_theme_falls_back(tmp_output_dir: Path):
    from app.tools.document_generator import create_pptx
    # Bogus theme should fall back to modern-dark — function still produces a valid file.
    result = create_pptx(
        title="Test",
        slides=[{"title": "Slide", "body": "Hello"}],
        style="solar-flare-explosion",
    )
    assert result["success"] is True


# ── helper ─────────────────────────────────────────────────────────────────

def _slide_text(slide) -> str:
    """Concatenate all visible text on a slide for assertions."""
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                parts.append(run.text)
    return " | ".join(p for p in parts if p)
